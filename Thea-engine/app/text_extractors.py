from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Tuple

from app.ocr import extract_text_from_image_path

try:
    from docx import Document
except Exception:
    Document = None

try:
    import openpyxl
    from openpyxl.utils import get_column_letter
except Exception:
    openpyxl = None
    get_column_letter = None

try:
    import xlrd
except Exception:
    xlrd = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import pytesseract
    from PIL import Image
except Exception:
    pytesseract = None
    Image = None


def extract_text_from_txt(file_path: Path) -> List[Tuple[int, str, Dict]]:
    raw_bytes = file_path.read_bytes()
    try:
        text = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = raw_bytes.decode("latin-1", errors="ignore")
    return [(1, text, {"type": "txt"})]


def extract_text_from_docx(file_path: Path) -> List[Tuple[int, str, Dict]]:
    if Document is None:
        raise Exception("python-docx not installed")
    doc = Document(str(file_path))
    pages: List[Tuple[int, str, Dict]] = []
    heading_stack: List[str] = []
    paragraph_index = 0

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name if paragraph.style else ""
        if style_name.lower().startswith("heading"):
            try:
                level = int(style_name.split()[-1])
            except Exception:
                level = 1
            heading_stack = heading_stack[: max(level - 1, 0)]
            heading_stack.append(text)
            continue

        paragraph_index += 1
        pages.append(
            (
                paragraph_index,
                text,
                {
                    "type": "docx",
                    "headingPath": " > ".join(heading_stack),
                    "paragraphIndex": paragraph_index,
                },
            )
        )

    if not pages:
        pages.append((1, "", {"type": "docx"}))
    return pages


def _sheet_to_text(rows: List[List[str]], sheet_name: str) -> str:
    lines = [f"Sheet: {sheet_name}"]
    for row in rows:
        if not row:
            continue
        lines.append(" | ".join(cell if cell is not None else "" for cell in row))
    return "\n".join(lines)


def extract_text_from_xlsx(file_path: Path) -> List[Tuple[int, str, Dict]]:
    if openpyxl is None:
        raise Exception("openpyxl not installed")
    workbook = openpyxl.load_workbook(str(file_path), data_only=True)
    pages: List[Tuple[int, str, Dict]] = []
    page_num = 0

    for sheet in workbook.worksheets:
        page_num += 1
        max_row = sheet.max_row or 0
        max_col = sheet.max_column or 0
        rows: List[List[str]] = []
        for row in sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
            rows.append([str(cell.value) if cell.value is not None else "" for cell in row])

        cell_range = "A1"
        if max_row > 0 and max_col > 0 and get_column_letter is not None:
            cell_range = f"A1:{get_column_letter(max_col)}{max_row}"

        pages.append(
            (
                page_num,
                _sheet_to_text(rows, sheet.title),
                {"type": "xlsx", "sheetName": sheet.title, "cellRange": cell_range},
            )
        )

    return pages or [(1, "", {"type": "xlsx"})]


def extract_text_from_xls(file_path: Path) -> List[Tuple[int, str, Dict]]:
    if xlrd is None:
        raise Exception("xlrd not installed")
    workbook = xlrd.open_workbook(str(file_path))
    pages: List[Tuple[int, str, Dict]] = []
    page_num = 0

    for sheet in workbook.sheets():
        page_num += 1
        rows: List[List[str]] = []
        for row_idx in range(sheet.nrows):
            row = sheet.row_values(row_idx)
            rows.append([str(cell) if cell is not None else "" for cell in row])

        cell_range = f"A1:{sheet.ncols}{sheet.nrows}" if sheet.nrows and sheet.ncols else "A1"
        pages.append(
            (
                page_num,
                _sheet_to_text(rows, sheet.name),
                {"type": "xls", "sheetName": sheet.name, "cellRange": cell_range},
            )
        )

    return pages or [(1, "", {"type": "xls"})]


def extract_text_from_pptx(file_path: Path) -> List[Tuple[int, str, Dict]]:
    if Presentation is None:
        raise Exception("python-pptx not installed")
    presentation = Presentation(str(file_path))
    pages: List[Tuple[int, str, Dict]] = []

    for index, slide in enumerate(presentation.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text or ""
            if notes_text.strip():
                texts.append(f"Notes: {notes_text.strip()}")
        page_text = "\n".join(texts).strip()
        pages.append((index, page_text, {"type": "pptx", "slideNumber": index}))

    return pages or [(1, "", {"type": "pptx"})]


def extract_text_from_image(file_path: Path) -> List[Tuple[int, str, Dict]]:
    if Image is None:
        raise Exception("Pillow not installed")
    image = Image.open(file_path)
    text = extract_text_from_image_path(file_path)
    bbox = None
    if pytesseract is not None:
        try:
            data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
            xs = [x for x in data.get("left", []) if x is not None]
            ys = [y for y in data.get("top", []) if y is not None]
            ws = [w for w in data.get("width", []) if w is not None]
            hs = [h for h in data.get("height", []) if h is not None]
            if xs and ys and ws and hs:
                x1 = min(xs)
                y1 = min(ys)
                x2 = max(x + w for x, w in zip(xs, ws))
                y2 = max(y + h for y, h in zip(ys, hs))
                bbox = [x1, y1, x2, y2]
        except Exception:
            bbox = None
    return [(1, text, {"type": "image", "bbox": bbox})]
